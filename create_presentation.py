"""
Generates Task_02 PowerPoint presentation.
Run: pip install python-pptx && python create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(1))
        tb2.text_frame.paragraphs[0].text = subtitle
        tb2.text_frame.paragraphs[0].font.size = Pt(24)
    return slide

def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].font.bold = True
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.space_after = Pt(6)
    return slide

def add_two_column_slide(title, left, right):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(32)
    tb.text_frame.paragraphs[0].font.bold = True
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6.3), Inches(5.5))
    for i, b in enumerate(left):
        p = left_box.text_frame.paragraphs[0] if i == 0 else left_box.text_frame.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.space_after = Pt(4)
    for i, b in enumerate(right):
        p = right_box.text_frame.paragraphs[0] if i == 0 else right_box.text_frame.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.space_after = Pt(4)
    return slide

add_title_slide("Task_02: Database Discovery on EC2", "Discover databases across AWS accounts • Expose via API • Proof-of-concept")
add_content_slide("Confidentiality Notice", [
    "This presentation is confidential and proprietary to AIRBUS.",
    "Unauthorized distribution, disclosure, or use is strictly prohibited.",
])

add_content_slide("Problem Statement", [
    "Discover all databases installed on EC2 instances across multiple AWS accounts",
    "Capture: database type, version, database sizing, system sizing",
    "Expose the output via an API",
    "Constraints: Read-only, no SSH, no production risk",
])

add_content_slide("High-Level Architecture", [
    "Management (hub) account: Discovery Lambda, DynamoDB, API Gateway",
    "Spoke accounts: EC2 instances with SSM agent",
    "Flow: EventBridge → Discovery Lambda → AssumeRole → SSM Run Command → Parse → DynamoDB",
    "API: API Gateway + Lambda reads DynamoDB",
])

add_two_column_slide("Who Does What", [
    "Intern A: Discovery",
    "• discovery_handler.py",
    "• discovery_python.py",
    "• ssm-document.json",
    "• DynamoDB schema",
    "• Speaks: architecture, SSM, DynamoDB",
    "",
    "Integration: Schema defined by A, consumed by B",
], [
    "Intern B: IAM & API",
    "• iam/*.json",
    "• api_handler.py",
    "• api-gateway-config.md",
    "• EXECUTION_GUIDE.md",
    "• Speaks: intro, API, disclaimers",
    "",
    "IAM defined by B, used by A",
])

add_content_slide("Key Components", [
    "IAM: Cross-account assume role, least privilege, no credentials exposed",
    "SSM: Runs discovery script on EC2 (no SSH); detects MySQL, PostgreSQL, MongoDB",
    "Discovery script: pgrep, version commands, du, /proc — all read-only",
    "DynamoDB: Stores account_id, instance_id, engine, version, sizing",
    "API: /health, /accounts, /accounts/{id}/instances, /databases",
])

add_content_slide("Edge Cases Handled", [
    "EC2 without SSM → Skipped, logged",
    "Database installed but not running → status: installed",
    "Multiple DBs on one instance → Separate records per DB",
    "Permission failure → discovery_status: failed, error stored",
    "Lambda timeout → Previous data remains; no partial write",
    "Unknown engine → Not detected; instance shows db_id: none if no known DBs",
])

add_content_slide("Security & Compliance", [
    "No SSH; IAM-based access via SSM",
    "Read-only: No DDL, DML, config writes",
    "CloudTrail: All AssumeRole, SendCommand, DynamoDB logged",
    "Least privilege: Scoped permissions per role",
    "What we do NOT do: Connect to DBs, modify files, access application data",
])

add_content_slide("Demo Flow", [
    "1. Architecture (Intern A)",
    "2. Trigger Discovery Lambda (Intern A)",
    "3. Show SSM command history (Intern A)",
    "4. Show DynamoDB items (Intern A)",
    "5. API calls: /health, /accounts, /databases (Intern B)",
    "6. Failure handling (optional)",
    "7. Limitations & disclaimers (Intern B)",
])

add_content_slide("Limitations (Explicit)", [
    "Only SSM-managed instances",
    "Only MySQL, PostgreSQL, MongoDB",
    "Batch only (not real-time)",
    "Single region",
    "No RDS/Aurora; no container discovery",
    "No API authentication in POC",
])

add_content_slide("What We Would Add for Production", [
    "API auth (IAM or API key)",
    "DynamoDB GSI + pagination",
    "Retries / DLQ for DynamoDB",
    "Multi-region support",
    "Unit tests, IaC (Terraform/CloudFormation)",
])

add_title_slide("Thank You", "Questions?")

prs.save("TASK_02_PRESENTATION.pptx")
print("Created: TASK_02_PRESENTATION.pptx")
